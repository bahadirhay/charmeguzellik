from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(22)

    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Akilli Randevu Yonetim Sistemi",
        "A'dan Z'ye isleyis, teyit/iptal akisi ve operasyonel izleme",
    )

    slides = [
        (
            "Proje Ozeti",
            [
                "Amac: Randevu surecini uçtan uca dijitallestirmek",
                "Hedef: Cakişma, gecikme ve iletisim hatalarini azaltmak",
                "Kapsam: Talep, onay, teyit, iptal, operasyon gunu, bildirim",
            ],
            "Bu sistem sadece randevu almak degil; operasyonu yonetmek icin tasarlandi.",
        ),
        (
            "Ana Kanallar",
            [
                "Web musteri formu (site uzerinden talep)",
                "Admin panel (yetkili/personel randevu girisi)",
                "Musteri aksiyon linki (teyit, iptal, takvim guncelle)",
            ],
            "Uc farkli giris noktasi tek veri modeline yazar, tutarlilik buradan gelir.",
        ),
        (
            "Durum Yasam Dongusu",
            [
                "pending -> approved -> confirmed",
                "cancelled / rejected",
                "checked_in / no_show",
                "Eski akis uyumlulugu: cancel_request",
            ],
            "Gercek operasyonu yansitan durumlar sayesinde raporlama netlesiyor.",
        ),
        (
            "Personel ve Yetki",
            [
                "Hizmet-personel eslesmesi Personel Planlama ile yonetilir",
                "Self yetkili kullanici sadece kendi hizmetlerini gorur",
                "Self kullanici sadece kendi adina randevu acabilir",
                "Kurallar backend'de de zorunludur (UI bypass edilmez)",
            ],
            "Yetki ihlalleri API seviyesinde engellenir; guvenli tasarim.",
        ),
        (
            "Musteri Teyit ve Iptal Akisi",
            [
                "Onayli randevuya musteriye link gonderilir",
                "Teyit ediyorum -> confirmed",
                "Randevumu iptal et -> cancelled",
                "Iki aksiyonda da Telegram + musteri bilgilendirmesi",
            ],
            "Musteri tek linkten sureci yonetir, ekip anlik haberdar olur.",
        ),
        (
            "1 Gun Once Hatirlatma",
            [
                "Yaklasik 24 saat sonraki approved randevular secilir",
                "Musteriye teyit/iptal linkli hatirlatma e-postasi gonderilir",
                "Telegram'da 'hatirlatma gonderildi' bilgisi paylasilir",
                "Tekrar gonderimi onlemek icin not izi dusulur",
            ],
            "No-show oranini dusuren kritik otomasyon adimi.",
        ),
        (
            "Bildirim ve Dayaniklilik",
            [
                "Telegram: yeni talep, teyit, iptal, guncelleme",
                "E-posta: durum bazli bilgilendirme",
                "SMTP kota asiminda fallback stratejisi",
                "UI teknik hata yerine aksiyon odakli yonlendirir",
            ],
            "Bildirim kanali aksasa bile operasyon devam eder.",
        ),
        (
            "Operasyon Panelleri",
            [
                "Randevu ve CRM listelerinde personel gorunurlugu",
                "Iptal ettiklerim: kim iptal etti bilgisi",
                "Operasyon gecmisi: geldi / gelmedi",
                "Takvimde teyitli randevu etiketi",
            ],
            "Saha ekibi tek ekrandan karar alabilecek netlige sahip olur.",
        ),
        (
            "Is Degeri",
            [
                "Cakişma ve yanlis randevu riski azalir",
                "Musteri teyit orani artar",
                "No-show ve iletisim gecikmeleri azalir",
                "Ekip ici koordinasyon ve izlenebilirlik artar",
            ],
            "Hem musteri deneyimi hem operasyon verimliligi birlikte iyilesir.",
        ),
        (
            "Sonraki Faz",
            [
                "Appointment event log tablosu",
                "Waitlist ve iptalde otomatik teklif",
                "Oda/cihaz kapasite modeli",
                "Dashboard KPI: teyit/no-show/iptal trendleri",
            ],
            "Sistemi tam olceklenebilir bir randevu platformuna tasiyoruz.",
        ),
    ]

    for title, bullets, notes in slides:
        add_bullet_slide(prs, title, bullets, notes)

    output = "Reservation-System-Sunum.pptx"
    prs.save(output)
    print(output)


if __name__ == "__main__":
    main()
